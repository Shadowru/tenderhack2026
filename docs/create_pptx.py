"""Generate a modern business presentation for TenderHack 2026."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
RED = RGBColor(0xE5, 0x39, 0x35)
RED_DARK = RGBColor(0xC6, 0x28, 0x28)
BLUE = RGBColor(0x19, 0x76, 0xD2)
DARK = RGBColor(0x33, 0x40, 0x59)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_GREEN = RGBColor(0x43, 0xA0, 0x47)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16, color=DARK, spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, line_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = line_color or color
        p.font.bold = is_bold
        p.space_after = Pt(font_size * (spacing - 1) * 2)
    return txBox


def arrow(s, x1, y1, x2, y2):
    """Simple line as thin rectangle."""
    thickness = Pt(3)
    if abs(x2 - x1) > abs(y2 - y1):
        left = min(x1, x2)
        w = abs(x2 - x1) or Pt(3)
        shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y1, w, thickness)
    else:
        top = min(y1, y2)
        h = abs(y2 - y1) or Pt(3)
        shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, top, thickness, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
    shape.line.fill.background()


def slide_header(slide, number, title):
    """Add consistent header bar to slide — projector-friendly bold style."""
    add_rect(slide, 0, 0, W, Inches(0.12), RED)
    add_text(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
             f"{number:02d}", font_size=18, color=RED, bold=True)
    add_text(slide, Inches(1.4), Inches(0.22), Inches(10), Inches(0.6),
             title, font_size=30, color=DARK, bold=True)
    # Thick line under header
    add_rect(slide, Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.04), RGBColor(0xE0, 0xE0, 0xE0))


def add_stat_card(slide, left, top, value, label, accent=RED):
    card = add_rounded_rect(slide, left, top, Inches(2.4), Inches(1.5), WHITE, RGBColor(0xD0, 0xD0, 0xD0))
    # Thick accent top bar
    add_rect(slide, left + Inches(0.02), top + Inches(0.02), Inches(2.36), Inches(0.1), accent)
    add_text(slide, left + Inches(0.3), top + Inches(0.3), Inches(1.8), Inches(0.6),
             value, font_size=36, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), top + Inches(0.95), Inches(2.0), Inches(0.4),
             label, font_size=13, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 1: Title
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)
add_rect(slide, 0, Inches(0), W, Inches(0.12), RED)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
         "ПЕРСОНАЛИЗИРОВАННЫЙ", font_size=44, color=WHITE, bold=True)
add_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1),
         "УМНЫЙ ПОИСК ПРОДУКЦИИ", font_size=44, color=RED, bold=True)

add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
         "Портал поставщиков — zakupki.mos.ru", font_size=20, color=LIGHT_GRAY)

add_rect(slide, Inches(1.5), Inches(4.5), Inches(3), Inches(0.04), RED)

add_text(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
         "TenderHack 2026", font_size=18, color=LIGHT_GRAY)
add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
         "tenderhack.extra.moscow", font_size=16, color=BLUE)

# ================================================================
# SLIDE 2: Problem
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, 1, "ПРОБЛЕМА")

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
         "537 000 товаров. 4 190 организаций. Одна поисковая строка.",
         font_size=22, color=DARK, bold=True)

# Three scenario cards
scenarios = [
    ("Больница", "Ищет «бумага»", "Нужна бумага для ЭКГ,\nтермобумага, офисная", RGBColor(0xE3, 0xF2, 0xFD)),
    ("Школа", "Ищет «бумага»", "Нужны бланки для грамот,\nсертификат-бумага", RGBColor(0xFB, 0xE9, 0xE7)),
    ("Колледж", "Ищет «бумага»", "Нужна одноразовая посуда\nиз ламинированной бумаги", RGBColor(0xE8, 0xF5, 0xE9)),
]

for i, (who, query, need, bg) in enumerate(scenarios):
    left = Inches(0.8 + i * 4.0)
    card = add_rounded_rect(slide, left, Inches(2.2), Inches(3.6), Inches(2.5), bg, RGBColor(0xE0, 0xE0, 0xE0))
    add_text(slide, left + Inches(0.3), Inches(2.4), Inches(3.0), Inches(0.4),
             who, font_size=20, color=DARK, bold=True)
    add_text(slide, left + Inches(0.3), Inches(2.85), Inches(3.0), Inches(0.3),
             query, font_size=14, color=GRAY)
    add_text(slide, left + Inches(0.3), Inches(3.4), Inches(3.0), Inches(1.0),
             need, font_size=13, color=DARK)

# Arrow + result
add_text(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
         "Сегодня: все получают одинаковый результат — SvetoCopy A4",
         font_size=16, color=RED, bold=True)

add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
         "Наша задача: каждый видит то, что нужно именно ему",
         font_size=18, color=ACCENT_GREEN, bold=True)

# ================================================================
# SLIDE 3: Solution — 4 levels
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, 2, "РЕШЕНИЕ — 4 УРОВНЯ ИНТЕЛЛЕКТА")

levels = [
    ("01", "Морфологический\nпоиск", "Русская морфология\nSubject extraction\n«ручка» ≠ «тачка с ручкой»", BLUE),
    ("02", "Опечатки\nи синонимы", "490K словарь\nEN→RU раскладка\n123 правила синонимов", RGBColor(0x7B, 0x1F, 0xA2)),
    ("03", "Персонализация", "Cold start из 2M контрактов\n4 190 профилей организаций\nDynamic re-ranking", RED),
    ("04", "AI-расширение\n(RAG)", "Qwen 2.5 7B (Ollama)\nES категории → LLM\nНечёткие запросы", ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(levels):
    left = Inches(0.6 + i * 3.1)
    top = Inches(1.3)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.9), top, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, left + Inches(0.1), top + Inches(0.9), Inches(2.5), Inches(0.8),
             title, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # Description card
    card = add_rounded_rect(slide, left, top + Inches(1.7), Inches(2.8), Inches(2.0), BG_GRAY)
    add_text(slide, left + Inches(0.2), top + Inches(1.9), Inches(2.4), Inches(1.6),
             desc, font_size=12, color=GRAY)

# Bottom stat
add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
         "Качество поиска: 93% — релевантный результат в топ-5 (134 тестовых запроса)",
         font_size=16, color=DARK, bold=True)

# ================================================================
# SLIDE 4: Process Flow — clean horizontal diagram
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_GRAY)
slide_header(slide, 3, "ПАЙПЛАЙН ОБРАБОТКИ ЗАПРОСА")

# Clean 3-row flow diagram with large readable blocks

# Row positions
ROW1_Y = Inches(1.2)   # Input processing
ROW2_Y = Inches(3.0)   # Search engine
ROW3_Y = Inches(4.8)   # Data layer

BOX_H = Inches(1.2)
BOX_W = Inches(2.0)
ARROW_COLOR = RGBColor(0xBB, 0xBB, 0xBB)

def flow_box(slide, left, top, title, subtitle, color, w=BOX_W):
    """Large readable process box."""
    shape = add_rounded_rect(slide, left, top, w, BOX_H, WHITE, color)
    # Color left bar
    add_rect(slide, left + Inches(0.03), top + Inches(0.1), Inches(0.08), BOX_H - Inches(0.2), color)
    add_text(slide, left + Inches(0.25), top + Inches(0.15), w - Inches(0.4), Inches(0.4),
             title, font_size=13, color=DARK, bold=True)
    add_text(slide, left + Inches(0.25), top + Inches(0.55), w - Inches(0.4), Inches(0.5),
             subtitle, font_size=10, color=GRAY)
    return shape

def flow_arrow_h(slide, x, y, length):
    """Horizontal arrow (thick rectangle)."""
    add_rect(slide, x, y - Pt(2), length, Pt(4), ARROW_COLOR)
    # Arrowhead triangle
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
        x + length - Inches(0.12), y - Inches(0.08), Inches(0.16), Inches(0.16))
    tri.fill.solid()
    tri.fill.fore_color.rgb = ARROW_COLOR
    tri.line.fill.background()
    tri.rotation = 90.0

def flow_arrow_v(slide, x, y, length):
    """Vertical arrow (thick rectangle)."""
    add_rect(slide, x - Pt(2), y, Pt(4), length, ARROW_COLOR)

# --- ROW 1: Input processing ---
add_text(slide, Inches(0.5), ROW1_Y - Inches(0.35), Inches(3), Inches(0.3),
         "ОБРАБОТКА ВВОДА", font_size=11, color=RED, bold=True)

flow_box(slide, Inches(0.5), ROW1_Y, "Ввод запроса", "Строка поиска\nАвтокомплит", BLUE)
flow_arrow_h(slide, Inches(2.55), ROW1_Y + BOX_H/2, Inches(0.35))
flow_box(slide, Inches(3.0), ROW1_Y, "Нормализация", "Раскладка EN→RU\nТранслит lat→кир", BLUE)
flow_arrow_h(slide, Inches(5.05), ROW1_Y + BOX_H/2, Inches(0.35))
flow_box(slide, Inches(5.5), ROW1_Y, "Spell Check", "490K словарь\nДетекция отрицания", BLUE)
flow_arrow_h(slide, Inches(7.55), ROW1_Y + BOX_H/2, Inches(0.35))
flow_box(slide, Inches(8.0), ROW1_Y, "Персонализация", "Cold start PG→Redis\nProduct + Category boosts", RGBColor(0x7B, 0x1F, 0xA2))
flow_arrow_h(slide, Inches(10.05), ROW1_Y + BOX_H/2, Inches(0.35))
flow_box(slide, Inches(10.5), ROW1_Y, "Результат", "Карточки + Формула\nAI-расширение (RAG)", ACCENT_GREEN)

# --- ROW 2: Search engine ---
add_text(slide, Inches(0.5), ROW2_Y - Inches(0.35), Inches(3), Inches(0.3),
         "ПОИСКОВЫЙ ДВИЖОК", font_size=11, color=RED, bold=True)

flow_box(slide, Inches(0.5), ROW2_Y, "Short mode", "Prefix subject.raw\nFuzzy subject (≤6 сим.)", RGBColor(0xF5, 0x70, 0x0C), w=Inches(2.3))
flow_box(slide, Inches(3.0), ROW2_Y, "Full mode", "Multi-match 4 стратегии\nPhrase + AND + Fuzzy", RGBColor(0xF5, 0x70, 0x0C), w=Inches(2.3))
flow_box(slide, Inches(5.6), ROW2_Y, "Subject Boost", "Exact +15 / Fuzzy +12\nSubject extraction", RED, w=Inches(2.3))
flow_box(slide, Inches(8.1), ROW2_Y, "Function Score", "BM25 × Popularity\n× Product × Category", RED, w=Inches(2.3))
flow_box(slide, Inches(10.6), ROW2_Y, "RAG (Ollama)", "ES категории → LLM\nQwen 2.5 7B", ACCENT_GREEN, w=Inches(2.3))

# --- ROW 3: Data layer ---
add_text(slide, Inches(0.5), ROW3_Y - Inches(0.35), Inches(3), Inches(0.3),
         "СЛОЙ ДАННЫХ", font_size=11, color=RED, bold=True)

flow_box(slide, Inches(0.5), ROW3_Y, "Elasticsearch", "537K docs, BM25\nSynonyms, Fuzzy", ACCENT_GREEN, w=Inches(2.8))
flow_box(slide, Inches(3.6), ROW3_Y, "Redis", "Profiles, Cache\nEvents, Cart", RGBColor(0xDC, 0x38, 0x2D), w=Inches(2.8))
flow_box(slide, Inches(6.7), ROW3_Y, "PostgreSQL", "2M Contracts\nMetrics, Sessions", BLUE, w=Inches(2.8))
flow_box(slide, Inches(9.8), ROW3_Y, "Ollama (CPU)", "Qwen 2.5 7B\nAI query expansion", RGBColor(0x7B, 0x1F, 0xA2), w=Inches(2.8))

# Vertical arrows between rows
for x_center in [Inches(1.5), Inches(4.15), Inches(6.75), Inches(9.25), Inches(11.65)]:
    flow_arrow_v(slide, x_center, ROW1_Y + BOX_H + Inches(0.05), Inches(0.45))

for x_center in [Inches(1.9), Inches(5.0), Inches(8.1), Inches(11.2)]:
    flow_arrow_v(slide, x_center, ROW2_Y + BOX_H + Inches(0.05), Inches(0.45))

# (old BPMN code removed)
# Placeholder to find next slide marker

# SLIDE 5: Personalization Demo
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, 4, "ПЕРСОНАЛИЗАЦИЯ В ДЕЙСТВИИ")

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         'Один запрос «бумага» — три разных результата:',
         font_size=18, color=DARK, bold=True)

demos = [
    ("Школа", "5003021368", ["Дизайн-бумага Металлик", "Сертификат-бумага Attache", "Бумага для творчества"], "Бланки для грамот\nи сертификатов", RGBColor(0xE3, 0xF2, 0xFD), BLUE),
    ("Колледж", "7716237684", ["Тарелка из ламинированной бумаги", "Стакан бумажный одноразовый", "Ролик подачи бумаги"], "Столовая + расходники\nдля принтеров", RGBColor(0xFB, 0xE9, 0xE7), RED),
    ("Больница", "7734091519", ["SvetoCopy Classic А4", "SvetoCopy (А4, 80 г/кв.м)", "ОФИСМАГ СТАНДАРТ А4"], "Стандартная\nофисная бумага", RGBColor(0xE8, 0xF5, 0xE9), ACCENT_GREEN),
]

for i, (who, inn, results, why, bg, accent) in enumerate(demos):
    left = Inches(0.6 + i * 4.1)
    top = Inches(1.8)

    card = add_rounded_rect(slide, left, top, Inches(3.8), Inches(4.2), bg, accent)

    add_text(slide, left + Inches(0.3), top + Inches(0.2), Inches(3.2), Inches(0.4),
             who, font_size=20, color=accent, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(0.6), Inches(3.2), Inches(0.3),
             f"ИНН: {inn}", font_size=10, color=GRAY)

    add_text(slide, left + Inches(0.3), top + Inches(1.0), Inches(3.2), Inches(0.3),
             "Результаты:", font_size=11, color=DARK, bold=True)

    for j, r in enumerate(results):
        add_text(slide, left + Inches(0.3), top + Inches(1.35 + j * 0.4), Inches(3.2), Inches(0.35),
                 f"  {j+1}. {r}", font_size=11, color=DARK)

    add_rect(slide, left + Inches(0.3), top + Inches(2.8), Inches(3.2), Inches(0.02), accent)

    add_text(slide, left + Inches(0.3), top + Inches(3.0), Inches(3.2), Inches(0.3),
             "Почему:", font_size=11, color=accent, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(3.3), Inches(3.2), Inches(0.7),
             why, font_size=11, color=GRAY)

# Bottom
add_text(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(0.4),
         "Персонализация подсказок тоже работает — автокомплит учитывает профиль организации",
         font_size=13, color=GRAY)

# ================================================================
# SLIDE 6: Key Metrics
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, 5, "КЛЮЧЕВЫЕ ПОКАЗАТЕЛИ")

# Stat cards — top row
add_stat_card(slide, Inches(0.6), Inches(1.2), "537 314", "Товаров проиндексировано", RED)
add_stat_card(slide, Inches(3.3), Inches(1.2), "2 009 457", "Контрактов проанализировано", BLUE)
add_stat_card(slide, Inches(6.0), Inches(1.2), "4 190", "Организаций с профилями", ACCENT_GREEN)
add_stat_card(slide, Inches(8.7), Inches(1.2), "490 204", "Терминов в словаре", RGBColor(0x7B, 0x1F, 0xA2))
add_stat_card(slide, Inches(11.4), Inches(1.2), "93%", "Качество поиска (топ-5)", RED)

# Metrics table
add_text(slide, Inches(0.8), Inches(3.1), Inches(5), Inches(0.4),
         "Метрики качества поиска", font_size=18, color=DARK, bold=True)

metrics = [
    ("MRR", "Позиция первого релевантного", "> 0.5"),
    ("Precision@10", "Доля полезных в топ-10", "> 40%"),
    ("CTR", "Клики / поисковые запросы", "> 30%"),
    ("Session Success", "Сессий с кликом", "> 50%"),
    ("Bounce Rate", "Быстрый возврат с карточки", "< 30%"),
    ("Personalization Lift", "Улучшение CTR при персонализации", "> 10%"),
]

for i, (name, desc, target) in enumerate(metrics):
    y = Inches(3.6 + i * 0.5)
    bg_c = BG_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.45), bg_c)
    add_text(slide, Inches(1.0), y + Inches(0.05), Inches(2.5), Inches(0.35),
             name, font_size=12, color=DARK, bold=True)
    add_text(slide, Inches(3.5), y + Inches(0.05), Inches(6.0), Inches(0.35),
             desc, font_size=12, color=GRAY)
    add_text(slide, Inches(10.0), y + Inches(0.05), Inches(2.0), Inches(0.35),
             target, font_size=12, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 7: Architecture
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, 6, "АРХИТЕКТУРА СИСТЕМЫ")

# Service boxes
services = [
    ("React SPA\nNginx", Inches(1.0), Inches(1.5), BLUE, "Frontend\nUI, карточки, корзина"),
    ("FastAPI\nPython 3.12", Inches(4.5), Inches(1.5), RED, "Backend\nAPI, персонализация"),
    ("Elasticsearch\n8.17", Inches(1.0), Inches(3.8), ACCENT_GREEN, "537K docs\nBM25, fuzzy, synonyms"),
    ("Redis 7", Inches(4.5), Inches(3.8), RGBColor(0xDC, 0x38, 0x2D), "Profiles, cache\nevents, cart"),
    ("PostgreSQL\n16", Inches(8.0), Inches(3.8), BLUE, "2M contracts\nmetrics, sessions"),
    ("Ollama\nQwen 2.5 7B", Inches(8.0), Inches(1.5), RGBColor(0x7B, 0x1F, 0xA2), "AI expansion\nRAG-поиск"),
]

for title, left, top, color, desc in services:
    card = add_rounded_rect(slide, left, top, Inches(2.8), Inches(1.6), WHITE, color)
    # Color bar on left
    add_rect(slide, left + Inches(0.02), top + Inches(0.1), Inches(0.08), Inches(1.4), color)

    add_text(slide, left + Inches(0.3), top + Inches(0.15), Inches(2.3), Inches(0.6),
             title, font_size=13, color=DARK, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(0.8), Inches(2.3), Inches(0.6),
             desc, font_size=10, color=GRAY)

# Caddy + SonarQube
add_text(slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.4),
         "Caddy (HTTPS, reverse proxy)  |  SonarQube (Quality Gate: OK)  |  Docker Compose  |  32 CPU, 125GB RAM",
         font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Arrows between services
arrow(slide, Inches(3.8), Inches(2.3), Inches(4.5), Inches(2.3))
arrow(slide, Inches(5.9), Inches(3.1), Inches(5.9), Inches(3.8))
arrow(slide, Inches(2.4), Inches(3.1), Inches(2.4), Inches(3.8))
arrow(slide, Inches(7.3), Inches(2.3), Inches(8.0), Inches(2.3))
arrow(slide, Inches(7.3), Inches(4.6), Inches(8.0), Inches(4.6))

# ================================================================
# SLIDE 8: Innovation highlights
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, 7, "КЛЮЧЕВЫЕ ИННОВАЦИИ")

innovations = [
    ("Subject Extraction", "Морфоанализ pymorphy3 выделяет ПРЕДМЕТ товара из названия. «Ручка» ≠ «тачка с ручкой». Падежи, сокращения, предлоги — всё учтено.", RED),
    ("Двухрежимный поиск", "Short (≤6 символов): prefix по subject.raw. Full: 4 параллельных стратегии (fuzzy + phrase×2 + AND). Оптимальная обработка каждого типа запроса.", BLUE),
    ("Cold Start", "4 190 организаций получают персонализацию с первого визита. 2M контрактов → предагрегированные buyer_category_weights → Redis за O(1).", ACCENT_GREEN),
    ("RAG через Ollama", "LLM (Qwen 2.5 7B) + реальные категории из ES индекса. «Медицинские перчатки» → нитриловые, хирургические, латексные. Не галлюцинирует.", RGBColor(0x7B, 0x1F, 0xA2)),
    ("Bounce Rate Tracking", "Открытие карточки < 3 сек → негативный сигнал quick_return. 5 типов событий с весами (view 2.0, click 3.0, cart 5.0). Temporal decay 0.95/день.", RGBColor(0xF5, 0x70, 0x0C)),
]

for i, (title, desc, color) in enumerate(innovations):
    top = Inches(1.2 + i * 1.15)
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top + Inches(0.05), Inches(0.45), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, Inches(1.5), top, Inches(2.5), Inches(0.35),
             title, font_size=15, color=DARK, bold=True)
    add_text(slide, Inches(4.2), top, Inches(8.5), Inches(0.9),
             desc, font_size=12, color=GRAY)

# ================================================================
# SLIDE 9: Ranking Signals — Positive & Negative
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, 8, "СИГНАЛЫ РАНЖИРОВАНИЯ")

# Two columns: positive (green) and negative (red)
add_text(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
         "ПОЗИТИВНЫЕ (повышают позицию)", font_size=20, color=ACCENT_GREEN, bold=True)
add_rect(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.05), ACCENT_GREEN)

positives = [
    ("Совпадение с предметом товара", "Subject extraction: поиск «ручка» → товары где ручка = предмет", "x5"),
    ("Точная фраза в названии", "match_phrase slop:1 — «100 листов» выше «500 листов»", "x20"),
    ("История закупок организации", "Cold start из контрактов: школа → бланки, больница → офисная", "x1.05"),
    ("Добавлен в корзину / избранное", "Event cart (вес 5.0) → мгновенный product boost", "x2.0"),
    ("Просмотр карточки > 3 сек", "Event view (вес 2.0) — заинтересовался, не bounce", "x1.5"),
    ("Популярность (кол-во контрактов)", "log2p(popularity) — мягкий множитель", "x1-3"),
]

for i, (signal, desc, boost) in enumerate(positives):
    y = Inches(1.95 + i * 0.7)
    add_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.6), BG_GRAY if i % 2 == 0 else WHITE)
    add_text(slide, Inches(1.0), y + Inches(0.05), Inches(3.5), Inches(0.25),
             signal, font_size=13, color=DARK, bold=True)
    add_text(slide, Inches(1.0), y + Inches(0.3), Inches(4.0), Inches(0.25),
             desc, font_size=10, color=GRAY)
    add_text(slide, Inches(5.2), y + Inches(0.1), Inches(1.0), Inches(0.35),
             boost, font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.RIGHT)

# Negative column
add_text(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.5),
         "НЕГАТИВНЫЕ (понижают / исключают)", font_size=20, color=RED, bold=True)
add_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.05), RED)

negatives = [
    ("Отрицание в запросе", "«не маслянных», «без латексных», «кроме» → must_not", "Исключ."),
    ("Bounce (< 3 сек на карточке)", "Event quick_return — закрыл карточку быстро", "-сигнал"),
    ("Слово — атрибут, не предмет", "«тачка с ручкой» → subject=тачка, ручка не предмет", "x0"),
    ("Temporal decay", "DECAY_FACTOR = 0.95 на каждый день давности события", "×0.95/д"),
    ("Историческая vs live-данные", "Контракты из PG дисконтируются ×0.7", "×0.7"),
]

for i, (signal, desc, effect) in enumerate(negatives):
    y = Inches(1.95 + i * 0.7)
    add_rect(slide, Inches(7.0), y, Inches(5.5), Inches(0.6), RGBColor(0xFE, 0xF0, 0xEF) if i % 2 == 0 else WHITE)
    add_text(slide, Inches(7.2), y + Inches(0.05), Inches(3.5), Inches(0.25),
             signal, font_size=13, color=DARK, bold=True)
    add_text(slide, Inches(7.2), y + Inches(0.3), Inches(4.0), Inches(0.25),
             desc, font_size=10, color=GRAY)
    add_text(slide, Inches(11.2), y + Inches(0.1), Inches(1.2), Inches(0.35),
             effect, font_size=14, color=RED, bold=True, alignment=PP_ALIGN.RIGHT)

# Bottom note
add_text(slide, Inches(0.8), Inches(6.4), Inches(12), Inches(0.4),
         "Формула: Score = BM25(strategies) x Popularity(log2p) x ProductBoost(1-2) x CategoryBoost(1-1.05) + SubjectExact(+15) + SubjectFuzzy(+12)",
         font_size=12, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 10: Thank you / CTA
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, W, Inches(0.12), RED)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
         "СПАСИБО", font_size=48, color=WHITE, bold=True)

add_rect(slide, Inches(1.5), Inches(2.6), Inches(3), Inches(0.04), RED)

# Links
links = [
    ("Демо:", "tenderhack.extra.moscow"),
    ("SonarQube:", "tenderhack.extra.moscow/sonar/"),
    ("GitHub:", "github.com/Shadowru/tenderhack2026"),
]

for i, (label, url) in enumerate(links):
    y = Inches(3.2 + i * 0.55)
    add_text(slide, Inches(1.5), y, Inches(2), Inches(0.4),
             label, font_size=16, color=LIGHT_GRAY)
    add_text(slide, Inches(3.5), y, Inches(8), Inches(0.4),
             url, font_size=16, color=BLUE, bold=True)

# Stats row at bottom
bottom_stats = [
    ("537K товаров", Inches(1.5)),
    ("2M контрактов", Inches(4.0)),
    ("4190 профилей", Inches(6.5)),
    ("93% качество", Inches(9.0)),
    ("30мс ответ", Inches(11.0)),
]

for text, left in bottom_stats:
    add_text(slide, left, Inches(5.5), Inches(2.5), Inches(0.4),
             text, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.4),
         "TenderHack 2026", font_size=14, color=LIGHT_GRAY)

# Save
output = "/home/zenith/tenderhack/docs/TenderHack2026_Presentation.pptx"
prs.save(output)
print(f"Created {output}")
